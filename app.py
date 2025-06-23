import os
from flask import Flask, render_template, request, send_file
from io import BytesIO
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import numpy as np
import wave
from PyPDF2 import PdfReader, PdfWriter
import mimetypes


app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ZW_SPACE = '\u200b'  # 0
ZW_NONJOIN = '\u200c'  # 1
END_MARK = '###FIN###'

# ===Txt Steganography ===

def cacher_texte(txt, secret):
    secret += END_MARK
    bits = ''.join(format(ord(c), '08b') for c in secret)
    zw_text = ''.join(ZW_SPACE if b == '0' else ZW_NONJOIN for b in bits)
    return txt + zw_text

def extraire_texte(txt):
    zw_seq = ''.join(c for c in txt if c in [ZW_SPACE, ZW_NONJOIN])
    if not zw_seq:
        return "Aucun texte caché trouvé."
    bits = ''.join('0' if c == ZW_SPACE else '1' for c in zw_seq)
    chars = []
    for i in range(0, len(bits), 8):
        byte = bits[i:i+8]
        if len(byte) == 8:
            chars.append(chr(int(byte, 2)))
    msg = ''.join(chars)
    fin = msg.find(END_MARK)
    return msg[:fin] if fin != -1 else "Aucun texte caché trouvé."

# === PNG Steganography ===

def cacher_texte_dans_png(img_bytes, texte):
    image = Image.open(BytesIO(img_bytes))
    img = image.convert("RGB")
    pixels = img.load()
    texte += END_MARK
    binaire = ''.join(format(ord(c), '08b') for c in texte)
    idx = 0
    for y in range(img.height):
        for x in range(img.width):
            if idx < len(binaire):
                r, g, b = pixels[x, y]
                r = (r & ~1) | int(binaire[idx])
                idx += 1
                if idx < len(binaire):
                    g = (g & ~1) | int(binaire[idx])
                    idx += 1
                if idx < len(binaire):
                    b = (b & ~1) | int(binaire[idx])
                    idx += 1
                pixels[x, y] = (r, g, b)
    out_bytes = BytesIO()
    img.save(out_bytes, format="PNG")
    out_bytes.seek(0)
    return out_bytes

def extraire_texte_dans_png(img_bytes):
    image = Image.open(BytesIO(img_bytes))
    img = image.convert('RGB')
    pixels = img.load()
    bits = ''
    for y in range(img.height):
        for x in range(img.width):
            r, g, b = pixels[x, y]
            bits += str(r & 1)
            bits += str(g & 1)
            bits += str(b & 1)
    chars = []
    for i in range(0, len(bits), 8):
        byte = bits[i:i+8]
        if len(byte) == 8:
            chars.append(chr(int(byte, 2)))
    message = ''.join(chars)
    fin = message.find(END_MARK)
    return message[:fin] if fin != -1 else "Aucun texte caché trouvé."

# === WAV Steganography ===

def cacher_texte_dans_wav(wav_bytes, texte):
    texte += END_MARK
    binaire = ''.join(format(ord(c), '08b') for c in texte)
    binaire += '0' * 8

    with wave.open(BytesIO(wav_bytes), 'rb') as wr:
        params = wr.getparams()
        frames = wr.readframes(wr.getnframes())
        audio = np.frombuffer(frames, dtype=np.int16)

    if len(binaire) > len(audio):
        raise ValueError("Le message est trop long pour ce fichier audio.")

    audio_mod = np.copy(audio)
    for idx, bit in enumerate(binaire):
        audio_mod[idx] = (audio_mod[idx] & ~1) | int(bit)

    out_bytes = BytesIO()
    with wave.open(out_bytes, 'wb') as ww:
        ww.setparams(params)
        ww.writeframes(audio_mod.tobytes())
    out_bytes.seek(0)
    return out_bytes

def extraire_texte_dans_wav(wav_bytes):
    with wave.open(BytesIO(wav_bytes), 'rb') as wr:
        frames = wr.readframes(wr.getnframes())
        audio = np.frombuffer(frames, dtype=np.int16)
    bits = ''.join(str(sample & 1) for sample in audio)
    chars = []
    for i in range(0, len(bits), 8):
        byte = bits[i:i+8]
        if len(byte) == 8:
            chars.append(chr(int(byte, 2)))
    message = ''.join(chars)
    fin = message.find(END_MARK)
    return message[:fin] if fin != -1 else "Aucun texte caché trouvé."

# === PDF Steganography (metadata) ===

def cacher_texte_dans_pdf(pdf_bytes, texte):
    texte += END_MARK
    reader = PdfReader(BytesIO(pdf_bytes))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.add_metadata({"/SteganoSecret": texte})
    output = BytesIO()
    writer.write(output)
    output.seek(0)
    return output

def extraire_texte_dans_pdf(pdf_bytes):
    reader = PdfReader(BytesIO(pdf_bytes))
    meta = reader.metadata
    secret = meta.get("/SteganoSecret")
    if secret and END_MARK in secret:
        return secret.split(END_MARK)[0]
    else:
        return "Aucun texte caché trouvé."

# === MP4 Steganography  ===

def inject_file_into_video(video_path, secret_path, output_path):
    with open(video_path, 'rb') as video, open(secret_path, 'rb') as secret:
        video_data = video.read()
        secret_data = secret.read()
    with open(output_path, 'wb') as out:
        out.write(video_data)
        out.write(b'\n--SECRET--\n')
        out.write(secret_data)

def extract_secret_from_video(file_path):
    with open(file_path, 'rb') as f:
        data = f.read()
    marker = b'\n--SECRET--\n'
    if marker in data:
        parts = data.split(marker)
        return parts[1]
    return None

@app.route('/', methods=['GET', 'POST'])
def index():
    message = ''
    if request.method == 'POST':
        action = request.form.get('action')
        file = request.files['texte']
        filename = file.filename

        if filename.endswith('.docx'):
            doc = Document(file)
            if action == 'cacher':
                texte_secret = request.form.get('texte', '')
                if doc.paragraphs:
                    doc.paragraphs[-1].text += cacher_texte('', texte_secret)
                else:
                    doc.add_paragraph(cacher_texte('', texte_secret))
                output = BytesIO()
                doc.save(output)
                output.seek(0)
                return send_file(output, as_attachment=True, download_name=f'stegano_{filename}', mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            elif action == 'extraire':
                hidden = ''.join(p.text for p in doc.paragraphs)
                message = "Texte caché extrait : " + extraire_texte(hidden)

        elif filename.endswith('.xlsx'):
            in_mem_file = BytesIO(file.read())
            wb = load_workbook(in_mem_file)
            ws = wb.active
            texte_secret = request.form.get('texte', '')
            if action == 'cacher':
                last_row = ws.max_row
                last_col = ws.max_column
                cell = ws.cell(row=last_row, column=last_col)
                cell.value = (cell.value or "") + cacher_texte('', texte_secret)
                output = BytesIO()
                wb.save(output)
                output.seek(0)
                return send_file(
                    output,
                    as_attachment=True,
                    download_name=f'stegano_{filename}',
                    mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                )
            elif action == 'extraire':
                all_text = ''
                for row in ws.iter_rows(values_only=True):
                    for val in row:
                        if val:
                            all_text += str(val)
                message = "Texte caché extrait : " + extraire_texte(all_text)

        elif filename.endswith('.pptx'):
            in_mem_file = BytesIO(file.read())
            prs = Presentation(in_mem_file)
            texte_secret = request.form.get('texte', '')
            if action == 'cacher':
                textbox_found = False
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape.text += cacher_texte('', texte_secret)
                        textbox_found = True
                        break
                if not textbox_found:
                    left = top = width = height = 100000
                    txBox = prs.slides[0].shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = cacher_texte('', texte_secret)
                output = BytesIO()
                prs.save(output)
                output.seek(0)
                return send_file(
                    output,
                    as_attachment=True,
                    download_name=f'stegano_{filename}',
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            elif action == 'extraire':
                all_text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            all_text += shape.text
                message = "Texte caché extrait : " + extraire_texte(all_text)

        elif filename.endswith('.png'):
            img_bytes = file.read()
            if action == 'cacher':
                texte_secret = request.form.get('texte', '')
                out_bytes = cacher_texte_dans_png(img_bytes, texte_secret)
                return send_file(out_bytes, as_attachment=True, download_name=f'stegano_{filename}', mimetype='image/png')
            elif action == 'extraire':
                message = "Texte caché extrait : " + extraire_texte_dans_png(img_bytes)

        elif filename.endswith('.wav'):
            wav_bytes = file.read()
            if action == 'cacher':
                texte_secret = request.form.get('texte', '')
                try:
                    out_bytes = cacher_texte_dans_wav(wav_bytes, texte_secret)
                except ValueError as e:
                    return render_template('index.html', message=str(e))
                return send_file(
                    out_bytes,
                    as_attachment=True,
                    download_name=f'stegano_{filename}',
                    mimetype='audio/wav'
                )
            elif action == 'extraire':
                message = "Texte caché extrait : " + extraire_texte_dans_wav(wav_bytes)

        elif filename.endswith('.pdf'):
            pdf_bytes = file.read()
            if action == 'cacher':
                texte_secret = request.form.get('texte', '')
                out_bytes = cacher_texte_dans_pdf(pdf_bytes, texte_secret)
                return send_file(
                    out_bytes,
                    as_attachment=True,
                    download_name=f'stegano_{filename}',
                    mimetype='application/pdf'
                )
            elif action == 'extraire':
                message = "Texte caché extrait : " + extraire_texte_dans_pdf(pdf_bytes)

        elif filename.endswith('.mp4'):
            video_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(video_path)

            if action == 'cacher':
                secret_file = request.files.get('secretfile')
                if not secret_file or secret_file.filename == '':
                    message = "Aucun fichier à cacher sélectionné."
                else:
                    secret_path = os.path.join(app.config['UPLOAD_FOLDER'], secret_file.filename)
                    secret_file.save(secret_path)
                    output_path = os.path.join(app.config['UPLOAD_FOLDER'], f'stegano_{filename}')
                    inject_file_into_video(video_path, secret_path, output_path)
                    return send_file(output_path, as_attachment=True, download_name=f'stegano_{filename}', mimetype='video/mp4')

            elif action == 'extraire':
                extracted_data = extract_secret_from_video(video_path)
                if extracted_data:
                    # Deviner le type MIME si possible
                    mime_type, _ = mimetypes.guess_type("fichier_extrait")
                    return send_file(BytesIO(extracted_data), as_attachment=True, download_name='fichier_extrait', mimetype=mime_type or 'application/octet-stream')
                else:
                    message = "Aucune donnée cachée trouvée dans la vidéo."

        else:
            content = file.read().decode('utf-8', errors='replace')
            if action == 'cacher':
                texte_secret = request.form.get('texte', '')
                new_txt = cacher_texte(content, texte_secret)
                output = BytesIO()
                output.write(new_txt.encode('utf-8'))
                output.seek(0)
                return send_file(output, as_attachment=True, download_name=f'stegano_{filename}', mimetype='text/plain')
            elif action == 'extraire':
                message = "Texte caché extrait : " + extraire_texte(content)

    return render_template('index.html', message=message)

if __name__ == '__main__':
    app.run(debug=True)
